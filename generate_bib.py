#!/usr/bin/env python
# coding: utf-8

# Copyright 2015 Nitin Pandya [nitin.pandya at gmail dot com]
# License: GPL <http://www.gnu.org/copyleft/gpl.html>
import csv
import sys, getopt
from pptx import Presentation

def main(argv):
    inputTemplate = ''
    inputCSV = ''
    try:
        opts, args = getopt.getopt(argv,"hi:c:",["ifile=","cfile="])
    except getopt.GetoptError:
        print 'generate_bib.py -i <input_ppt_template> -c <input_CSV>'
        sys.exit(2)
    for opt, arg in opts:
        if opt == '-h':
            print 'generate_bib.py -i <input_ppt_template> -c <input_CSV>'
            sys.exit()
        elif opt in ("-i", "--ifile"):
            inputTemplate = arg
        elif opt in ("-c", "-cfile"):
            inputCSV = arg

    print("Generating bibs using base powerpoint template : " + inputTemplate + "and csv file " + inputCSV) 

	# Open a presentation
    f = open(inputTemplate)
    prs = Presentation(f)
    f.close()

    # Open CSV file
    csv_in = open(inputCSV)
    participants = csv.DictReader(csv_in, delimiter=',')

	# Get the templated slide layout    
    slide_layout = prs.slide_layouts[1]

    # for every CSV entry 
    for line in participants:
        print(line["Bib"], line["Name"], line["Category"])
        # Add new slide
    	slide = prs.slides.add_slide(slide_layout)
    	for shape in slide.placeholders:     		#print('%d %s' % (shape.placeholder_format.idx, shape.name))
    		if shape.placeholder_format.idx == 11:
	    		shape.text = line["Category"]
	    	elif shape.placeholder_format.idx == 12:
	    		shape.text = line["Name"]
	    	elif shape.placeholder_format.idx == 13:
	    		shape.text = line["Bib"]

	prs.save('generated_bibs.pptx')

if __name__ == "__main__":
    main(sys.argv[1:])